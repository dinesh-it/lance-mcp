#!/usr/bin/env python3
"""
Document to PDF Converter

Converts various document formats to PDF using Python libraries.
Supports: docx, doc, pptx, ppt, xlsx, xls, txt, csv, rtf, html, md
"""

import os
import sys
import argparse
from pathlib import Path
from typing import List, Dict, Set
import logging

# Document conversion libraries
try:
    from docx2pdf import convert as docx_to_pdf
    DOCX2PDF_AVAILABLE = True
except ImportError:
    DOCX2PDF_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.units import inch
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import markdown
    from weasyprint import HTML, CSS
    WEASYPRINT_AVAILABLE = True
    MARKDOWN_AVAILABLE = True
except ImportError:
    WEASYPRINT_AVAILABLE = False
    try:
        import markdown
        MARKDOWN_AVAILABLE = True
    except ImportError:
        MARKDOWN_AVAILABLE = False

# Setup logging
logging.basicConfig(level=logging.INFO, format='%(levelname)s: %(message)s')
logger = logging.getLogger(__name__)

# Supported file extensions mapped to conversion functions
SUPPORTED_EXTENSIONS = {
    '.docx': 'convert_docx',
    '.doc': 'convert_doc',
    '.pptx': 'convert_pptx', 
    '.ppt': 'convert_ppt',
    '.xlsx': 'convert_xlsx',
    '.xls': 'convert_xls',
    '.txt': 'convert_text',
    '.csv': 'convert_csv',
    '.rtf': 'convert_rtf',
    '.html': 'convert_html',
    '.htm': 'convert_html',
    '.md': 'convert_markdown'
}

class DocumentConverter:
    def __init__(self, overwrite: bool = False, keep_originals: bool = True):
        self.overwrite = overwrite
        self.keep_originals = keep_originals
        self.success_count = 0
        self.fail_count = 0
        self.check_dependencies()
    
    def check_dependencies(self):
        """Check which conversion libraries are available"""
        missing_deps = []
        
        if not REPORTLAB_AVAILABLE:
            missing_deps.append("reportlab (for text/CSV conversion)")
        if not PANDAS_AVAILABLE:
            missing_deps.append("pandas openpyxl (for Excel files)")
        if not DOCX2PDF_AVAILABLE:
            missing_deps.append("docx2pdf (for Word files)")
        if not PPTX_AVAILABLE:
            missing_deps.append("python-pptx (for PowerPoint files)")
        if not WEASYPRINT_AVAILABLE:
            missing_deps.append("weasyprint (for HTML conversion)")
        if not MARKDOWN_AVAILABLE:
            missing_deps.append("markdown (for Markdown files)")
        
        if missing_deps:
            logger.warning("Missing optional dependencies:")
            for dep in missing_deps:
                logger.warning(f"  - {dep}")
            logger.warning("Install with: pip install -r requirements.txt")
    
    def convert_docx(self, input_file: Path, output_file: Path) -> bool:
        """Convert DOCX to PDF using docx2pdf"""
        if not DOCX2PDF_AVAILABLE:
            logger.error("docx2pdf not available. Install with: pip install docx2pdf")
            return False
        
        try:
            docx_to_pdf(str(input_file), str(output_file))
            return True
        except Exception as e:
            logger.error(f"Failed to convert {input_file}: {e}")
            return False
    
    def convert_doc(self, input_file: Path, output_file: Path) -> bool:
        """Convert DOC to PDF (requires LibreOffice or docx2pdf fallback)"""
        # docx2pdf can sometimes handle .doc files
        if DOCX2PDF_AVAILABLE:
            try:
                docx_to_pdf(str(input_file), str(output_file))
                return True
            except Exception as e:
                logger.error(f"Failed to convert {input_file}: {e}")
                return False
        
        logger.error("No converter available for .doc files. Consider converting to .docx first.")
        return False
    
    def convert_pptx(self, input_file: Path, output_file: Path) -> bool:
        """Convert PPTX to PDF using python-pptx and reportlab"""
        if not PPTX_AVAILABLE or not REPORTLAB_AVAILABLE:
            logger.error("python-pptx and reportlab required for PowerPoint conversion")
            return False
        
        try:
            # Load presentation
            prs = Presentation(str(input_file))
            
            # Create PDF document
            doc = SimpleDocTemplate(str(output_file), pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Add title
            title = Paragraph(f"<b>{input_file.name}</b>", styles['Title'])
            story.append(title)
            story.append(Spacer(1, 12))
            
            # Process each slide
            for i, slide in enumerate(prs.slides, 1):
                # Add slide header
                slide_header = Paragraph(f"<b>Slide {i}</b>", styles['Heading1'])
                story.append(slide_header)
                story.append(Spacer(1, 6))
                
                # Extract text from all shapes in the slide
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    # Join all text from the slide
                    combined_text = '\n\n'.join(slide_text)
                    # Add slide content
                    content_para = Paragraph(combined_text.replace('\n', '<br/>'), styles['Normal'])
                    story.append(content_para)
                else:
                    # No text content found
                    no_text = Paragraph("<i>No text content found in this slide</i>", styles['Italic'])
                    story.append(no_text)
                
                story.append(Spacer(1, 12))
            
            # Build PDF
            doc.build(story)
            return True
            
        except Exception as e:
            logger.error(f"Failed to convert {input_file}: {e}")
            return False
    
    def convert_ppt(self, input_file: Path, output_file: Path) -> bool:
        """Convert PPT to PDF"""
        return self.convert_pptx(input_file, output_file)  # Try same method
    
    def convert_xlsx(self, input_file: Path, output_file: Path) -> bool:
        """Convert XLSX to PDF using pandas and reportlab"""
        if not PANDAS_AVAILABLE or not REPORTLAB_AVAILABLE:
            logger.error("pandas and reportlab required for Excel conversion")
            return False
        
        try:
            # Read Excel file
            df = pd.read_excel(input_file, sheet_name=None)  # Read all sheets
            
            # Create PDF
            doc = SimpleDocTemplate(str(output_file), pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Add title
            title = Paragraph(f"<b>{input_file.name}</b>", styles['Title'])
            story.append(title)
            story.append(Spacer(1, 12))
            
            # Add each sheet
            for sheet_name, sheet_df in df.items():
                if len(df) > 1:  # Multiple sheets
                    sheet_title = Paragraph(f"<b>Sheet: {sheet_name}</b>", styles['Heading1'])
                    story.append(sheet_title)
                    story.append(Spacer(1, 12))
                
                # Convert DataFrame to string representation
                df_string = sheet_df.to_string(index=False)
                para = Paragraph(f"<pre>{df_string}</pre>", styles['Code'])
                story.append(para)
                story.append(Spacer(1, 12))
            
            doc.build(story)
            return True
            
        except Exception as e:
            logger.error(f"Failed to convert {input_file}: {e}")
            return False
    
    def convert_xls(self, input_file: Path, output_file: Path) -> bool:
        """Convert XLS to PDF"""
        return self.convert_xlsx(input_file, output_file)
    
    def convert_text(self, input_file: Path, output_file: Path) -> bool:
        """Convert text file to PDF using reportlab"""
        if not REPORTLAB_AVAILABLE:
            logger.error("reportlab required for text conversion")
            return False
        
        try:
            with open(input_file, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            doc = SimpleDocTemplate(str(output_file), pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Add title
            title = Paragraph(f"<b>{input_file.name}</b>", styles['Title'])
            story.append(title)
            story.append(Spacer(1, 12))
            
            # Add content (preserve formatting with <pre> tag)
            content_para = Paragraph(f"<pre>{content}</pre>", styles['Code'])
            story.append(content_para)
            
            doc.build(story)
            return True
            
        except Exception as e:
            logger.error(f"Failed to convert {input_file}: {e}")
            return False
    
    def convert_csv(self, input_file: Path, output_file: Path) -> bool:
        """Convert CSV to PDF using pandas and reportlab"""
        if not PANDAS_AVAILABLE or not REPORTLAB_AVAILABLE:
            logger.error("pandas and reportlab required for CSV conversion")
            return False
        
        try:
            # Read CSV
            df = pd.read_csv(input_file)
            
            # Create PDF
            doc = SimpleDocTemplate(str(output_file), pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Add title
            title = Paragraph(f"<b>{input_file.name}</b>", styles['Title'])
            story.append(title)
            story.append(Spacer(1, 12))
            
            # Convert DataFrame to string
            df_string = df.to_string(index=False)
            para = Paragraph(f"<pre>{df_string}</pre>", styles['Code'])
            story.append(para)
            
            doc.build(story)
            return True
            
        except Exception as e:
            logger.error(f"Failed to convert {input_file}: {e}")
            return False
    
    def convert_rtf(self, input_file: Path, output_file: Path) -> bool:
        """Convert RTF to PDF (basic text extraction)"""
        # RTF is complex, but we can try to extract basic text
        try:
            with open(input_file, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            # Basic RTF text extraction (very simple)
            # Remove RTF control words (this is very basic)
            import re
            text = re.sub(r'\\[a-z]+\d*\s?', '', content)
            text = re.sub(r'[{}]', '', text)
            text = text.strip()
            
            return self.convert_text_content(text, input_file.name, output_file)
            
        except Exception as e:
            logger.error(f"Failed to convert RTF {input_file}: {e}")
            return False
    
    def convert_html(self, input_file: Path, output_file: Path) -> bool:
        """Convert HTML to PDF using weasyprint"""
        if not WEASYPRINT_AVAILABLE:
            logger.error("weasyprint required for HTML conversion")
            return False
        
        try:
            HTML(filename=str(input_file)).write_pdf(str(output_file))
            return True
        except Exception as e:
            logger.error(f"Failed to convert {input_file}: {e}")
            return False
    
    def convert_markdown(self, input_file: Path, output_file: Path) -> bool:
        """Convert Markdown to PDF via HTML"""
        if not MARKDOWN_AVAILABLE:
            logger.error("markdown required for Markdown conversion")
            return False
        
        try:
            with open(input_file, 'r', encoding='utf-8') as f:
                md_content = f.read()
            
            # Convert to HTML
            html_content = markdown.markdown(md_content)
            
            if WEASYPRINT_AVAILABLE:
                # Use weasyprint for better formatting
                HTML(string=html_content).write_pdf(str(output_file))
            else:
                # Fallback to text conversion
                import re
                text = re.sub(r'<[^>]+>', '', html_content)  # Strip HTML tags
                return self.convert_text_content(text, input_file.name, output_file)
            
            return True
            
        except Exception as e:
            logger.error(f"Failed to convert {input_file}: {e}")
            return False
    
    def convert_text_content(self, content: str, filename: str, output_file: Path) -> bool:
        """Helper to convert text content to PDF"""
        if not REPORTLAB_AVAILABLE:
            return False
        
        try:
            doc = SimpleDocTemplate(str(output_file), pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            title = Paragraph(f"<b>{filename}</b>", styles['Title'])
            story.append(title)
            story.append(Spacer(1, 12))
            
            content_para = Paragraph(f"<pre>{content}</pre>", styles['Code'])
            story.append(content_para)
            
            doc.build(story)
            return True
        except Exception as e:
            logger.error(f"Failed to create PDF: {e}")
            return False
    
    def convert_document(self, input_file: Path, output_dir: Path) -> bool:
        """Convert a single document to PDF"""
        ext = input_file.suffix.lower()
        
        if ext not in SUPPORTED_EXTENSIONS:
            logger.warning(f"Unsupported file format: {ext}")
            return False
        
        output_file = output_dir / f"{input_file.stem}.pdf"
        
        # Check if output already exists
        if output_file.exists() and not self.overwrite:
            logger.info(f"Skipping {input_file.name} - PDF already exists")
            return True
        
        # Get conversion method
        method_name = SUPPORTED_EXTENSIONS[ext]
        method = getattr(self, method_name)
        
        logger.info(f"Converting: {input_file.name}")
        success = method(input_file, output_file)
        
        if success:
            self.success_count += 1
            # Remove original if requested
            if not self.keep_originals:
                try:
                    input_file.unlink()
                    logger.info(f"Removed original: {input_file.name}")
                except Exception as e:
                    logger.warning(f"Failed to remove original {input_file}: {e}")
        else:
            self.fail_count += 1
        
        return success
    
    def find_documents(self, input_dir: Path) -> List[Path]:
        """Find all supported documents in directory (recursive)"""
        documents = []
        
        for ext in SUPPORTED_EXTENSIONS.keys():
            documents.extend(input_dir.rglob(f"*{ext}"))
        
        return sorted(documents)
    
    def convert_all(self, input_dir: Path, output_dir: Path):
        """Convert all documents in input directory"""
        # Create output directory
        output_dir.mkdir(parents=True, exist_ok=True)
        
        # Find documents
        logger.info("Scanning for documents to convert...")
        documents = self.find_documents(input_dir)
        
        if not documents:
            logger.info("No supported documents found.")
            logger.info(f"Supported formats: {', '.join(SUPPORTED_EXTENSIONS.keys())}")
            return
        
        logger.info(f"Found {len(documents)} documents to convert")
        
        # Convert each document
        for doc in documents:
            self.convert_document(doc, output_dir)
        
        # Summary
        logger.info(f"\nConversion completed:")
        logger.info(f"  Successfully converted: {self.success_count}")
        logger.info(f"  Failed to convert: {self.fail_count}")
        logger.info(f"  Output directory: {output_dir}")

def main():
    parser = argparse.ArgumentParser(description="Convert documents to PDF format")
    parser.add_argument("--input", required=True, help="Input directory containing documents")
    parser.add_argument("--output", required=True, help="Output directory for PDF files")
    parser.add_argument("--overwrite", action="store_true", help="Overwrite existing PDF files")
    parser.add_argument("--keep-originals", action="store_true", default=True, 
                       help="Keep original files after conversion (default: True)")
    parser.add_argument("--remove-originals", action="store_true", 
                       help="Remove original files after successful conversion")
    
    args = parser.parse_args()
    
    # Handle the keep_originals logic
    keep_originals = args.keep_originals and not args.remove_originals
    
    input_dir = Path(args.input)
    output_dir = Path(args.output)
    
    if not input_dir.exists():
        logger.error(f"Input directory does not exist: {input_dir}")
        sys.exit(1)
    
    logger.info(f"INPUT DIRECTORY: {input_dir}")
    logger.info(f"OUTPUT DIRECTORY: {output_dir}")
    logger.info(f"OVERWRITE FLAG: {args.overwrite}")
    logger.info(f"KEEP ORIGINALS: {keep_originals}")
    
    converter = DocumentConverter(overwrite=args.overwrite, keep_originals=keep_originals)
    converter.convert_all(input_dir, output_dir)

if __name__ == "__main__":
    main()